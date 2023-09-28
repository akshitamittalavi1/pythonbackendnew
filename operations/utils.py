import openai
import fitz
import logging
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
import textwrap
from config import BLOB_CONN_STRING,BLOB_CONTAINER1, BLOB_CONTAINER2, OPENAI_API_TYPE, OPENAI_API_KEY, OPENAI_API_BASE,OPENAI_API_VERSION
from project_logs import MyLogger


# initialize openai API credentials and details
openai.api_key = OPENAI_API_KEY
openai.api_type = OPENAI_API_TYPE
openai.api_base = OPENAI_API_BASE
openai.api_version = OPENAI_API_VERSION

logger = MyLogger(filename="./projectlogs/utils_log.log")

class Utils:

    @logger.log_execution(logger)
    def get_embeddings_text(self, text):
        para = []
        embeddings = []

        # Initialize OpenAI API client
        try:
            for i in text:
                para.append(i)
                response = openai.Embedding.create(input=[i], engine="qna-embedding-model")
                embeddings.append(response['data'][0]['embedding'])
            logger.log(message=f"Paragraph and embedding returned Successfully", level=logging.INFO)
            print(f"Paragraph and embedding returned Successfully")
            return para, embeddings
        except Exception as e:
            # logger.log(message=f"{e}, function: get_embeddings_text", level=logging.ERROR)
            print(e)

    @logger.log_execution(logger)
    def create_doc(self, paragraphs, embeddings):
        documents = []
        try:
            for i in range(len(paragraphs)):
                documents.append({
                    "id": str(i),
                    "content": paragraphs[i],
                    "contentVector": embeddings[i]
                })
            print("Created Documents!")
            logger.log(message=f"Created Documents Successfully", level=logging.INFO)
            return documents
        except Exception as e:
            # logger.log(message=f"{e}, function: create_doc", level=logging.ERROR)
            print(e)

    @logger.log_execution(logger)
    def split_into_chunks(self, text):
        lines = text.split('\n')
        chunks = []
        current_chunk = ""
        try:
            for line in lines:
                line = line.strip()
                if line and line.isupper() and line.isupper() and len(line) > 5:
                    line = line.lower()
                    if current_chunk:
                        chunks.append(current_chunk)
                    current_chunk = line + '\n'
                else:
                    current_chunk += line + '\n'

            if current_chunk:
                chunks.append(current_chunk)
            # chunks=re.sub("[^ A-Za-z0-9]"," ",chunks)
            logger.log(message=f"Chunks Returned Successfully", level=logging.INFO)
            print("Chunks Returned Successfully")
            return chunks
        except Exception as e:
            # logger.log(message=f"{e}, function: split_into_chunks", level=logging.ERROR)
            print(e)

    @logger.log_execution(logger)
    def highlight_pdf_page(self, input_pdf, output_pdf, page_index_to_highlight):
#         """
#         Highlight a specific page in a PDF and save the result to a new PDF.

#         Args:
#             input_pdf (str): The path to the input PDF.
#             output_pdf (str): The path to save the highlighted PDF.
#             page_index_to_highlight (int): The index of the page to highlight (starting from 0).
#         """
        try:
            doc = fitz.open(input_pdf)

            if page_index_to_highlight < 0 or page_index_to_highlight >= len(doc):
                print(f"Invalid page index {page_index_to_highlight}. The PDF has {len(doc)} pages.")
                return

            page_data = doc[page_index_to_highlight]
            highlight_rectangle = page_data.rect
            highlight = page_data.add_highlight_annot(highlight_rectangle)

            doc.save(output_pdf)
            doc.close()

            print(f"Page {page_index_to_highlight} highlighted and saved to {output_pdf}")
        except Exception as e:
            print(f"An error occurred: {str(e)}")

# Example usage:
# highlight_pdf_page('HR Manual.pdf', 'highlighted_manual.pdf', 7)

    def calculate_tokens(self, page_data):
        """
        Calculate the token count based on the given page data.

        Args:
            page_data (str): The input text content of a page.

        Returns:
            float: The calculated token count, multiplied by 1.33.

        Raises:
            ValueError: If the input page_data is not a valid string.

        Example:
            page_data = "This is a sample page with some text."
            result = calculate_tokens(page_data)
            print(result)  # Output will be the calculated token count.
        """
        try:
            # Ensure that page_data is a valid string
            if not isinstance(page_data, str):
                raise ValueError("page_data must be a valid string")

            # Split the page data into tokens using spaces
            tokens = page_data.split()

            # Calculate the total number of tokens
            total_tokens = len(tokens)

            # Multiply the total tokens by 1.33
            final_value = total_tokens * 1.33

            return final_value
        except Exception as e:
            print(f"An error occurred: {str(e)}")
            return None  # Return None in case of an error


    
    def document_qa(self, query, language_type, user_id, question_type):
        try:
            if language_type=='ar':
                query=Translator.translate(query)

            if question_type=='docs':
                answer=OpenAIPipeline.generate_answer(query)
                return answer
            
            elif question_type=='bing':
                answer=Bing.single_answer(query)
                return answer
            document_path='To_Send\data\semantic-kernel-get-started-quick-start-guide.pdf'
            output_path='To_Send\temp'
            page_no=5
            self.highlight_pdf_page(document_path, output_path, page_no)

            
        except Exception as e:
            print(f'Exception occured {e}')
            
    def convert_ppt_to_pdf(self, pptx_path, pdf_path):
        # Load the PowerPoint presentation
        ppt = Presentation(pptx_path)

        # Create a PDF file
        c = canvas.Canvas(pdf_path, pagesize=letter)

        # Define the page width and height
        page_width, page_height = letter
        page_height -= 0.5 * inch  # leave some margin at the bottom

        # Iterate through slides and add them to the PDF
        for slide_number, slide in enumerate(ppt.slides):
            if slide_number != 0:
                c.showPage()  # Move to the next page
            
            c.drawString(inch, page_height - inch, f"Slide {slide_number + 1}:")

            y = page_height - 1.5 * inch  # Start below the slide number
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines = shape.text.split("\n")
                    for line in lines:
                        wrapped_lines = textwrap.wrap(line, width=60)  # Adjust width as needed
                        for wrapped_line in wrapped_lines:
                            if y - 0.5 * inch < 0:
                                c.showPage()  # Move to the next page
                                y = page_height - 1.5 * inch
                                c.drawString(inch, page_height - inch, f"Slide {slide_number + 1}:")
                            
                            c.drawString(inch * 1.5, y, wrapped_line)
                            y -= 0.5 * inch

        c.save()  # Save the PDF file

# Example usage:
# page_data = "This is a sample page with some text."
# result = calculate_tokens(page_data)
# if result is not None:
#     print(result)